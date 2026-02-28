from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from typing import Optional
import os
import asyncio
from pathlib import Path
import PyPDF2
import docx
from pptx import Presentation
import requests
import json
from gtts import gTTS
import uuid

# Import advanced universal analyzer
from advanced_universal_analyzer import analyze_code_advanced_universal

app = FastAPI(title="AI Document Reader Service", version="1.0.0")

# Request models
class FileProcessRequest(BaseModel):
    filePath: str
    fileName: str
    fileType: str

class ContentRequest(BaseModel):
    filePath: str
    contentType: str  # 'full' or 'summary'

class AudioRequest(BaseModel):
    text: str
    language: str = 'en'

# File processors
class DocumentProcessor:
    @staticmethod
    def extract_pdf_text(file_path: str) -> str:
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + " "
                
                # COMPREHENSIVE word fixing - ENHANCED VERSION
                import re
                
                # Step 1: Fix broken words with spaces in middle (SUPER ENHANCED)
                # Multiple passes for complex patterns
                
                # Pass 1: Fix 3-part broken words like "Det a iled" -> "Detailed"
                text = re.sub(r'([A-Z][a-z]{1,3})\s+([a-z]{1,3})\s+([a-z]{1,5})', r'\1\2\3', text)
                
                # Pass 2: Fix 2-part broken words like "T ask" -> "Task"  
                text = re.sub(r'([A-Z])\s+([a-z]{3,})', r'\1\2', text)
                
                # Pass 3: Fix common patterns like "Man a gement" -> "Management"
                text = re.sub(r'([A-Za-z]{3,})\s+([a-z]{1,5})([A-Za-z])', r'\1\2\3', text)
                text = re.sub(r'([A-Za-z]{3,})\s+([a-z]{1,5})\s+([A-Za-z])', r'\1\2 \3', text)
                
                # Pass 4: Aggressive specific fixes (case-insensitive)
                specific_fixes = [
                    ('Det a iled', 'Detailed'),
                    ('det a iled', 'detailed'),
                    ('T ask', 'Task'),
                    ('t ask', 'task'),
                    ('Br e a kdown', 'Breakdown'),
                    ('br e a kdown', 'breakdown'),
                    ('Man a gement', 'Management'),
                    ('man a gement', 'management'),
                    ('F ullstack', 'Fullstack'),
                    ('f ullstack', 'fullstack'),
                    ('R eact', 'React'),
                    ('r eact', 'react'),
                    ('V ite', 'Vite'),
                    ('v ite', 'vite'),
                    ('e xpert', 'expert'),
                    ('po wered', 'powered'),
                    ('Git Hub', 'GitHub'),
                    ('git hub', 'github'),
                    ('Builda', 'Build a'),
                    ('builda', 'build a'),
                    ('needa', 'need a'),
                    ('whois', 'who is'),
                    ('Systemwith', 'System with'),
                    ('systemwith', 'system with')
                ]
                
                for broken, fixed in specific_fixes:
                    text = text.replace(broken, fixed)
                
                # Step 2: Fix words stuck together with 'a': "BuildaFullstack" -> "Build a Fullstack"
                text = re.sub(r'([a-z])a([A-Z][a-z]+)', r'\1 a \2', text)
                text = re.sub(r'([a-z])a([a-z]{4,})', r'\1 a \2', text)
                
                # Step 3: Fix camelCase words: "needadeveloper" -> "need a developer"
                text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
                
                # Step 4: Fix words stuck together: "whoisane" -> "who is an e"
                text = re.sub(r'([a-z]{3,})is([a-z]{2,})', r'\1 is \2', text)
                text = re.sub(r'([a-z]{3,})be([a-z]{3,})', r'\1 be \2', text)
                text = re.sub(r'([a-z]{3,})the([a-z]{3,})', r'\1 the \2', text)
                text = re.sub(r'([a-z]{3,})to([a-z]{3,})', r'\1 to \2', text)
                text = re.sub(r'([a-z]{3,})for([a-z]{3,})', r'\1 for \2', text)
                text = re.sub(r'([a-z]{3,})with([a-z]{3,})', r'\1 with \2', text)
                text = re.sub(r'([a-z]{3,})and([a-z]{3,})', r'\1 and \2', text)
                
                # Step 5: Fix specific common broken patterns
                replacements = {
                    # Broken words with spaces
                    'Manag ement': 'Management',
                    'manag ement': 'management',
                    'Enquiry Manag ement': 'Enquiry Management',
                    'Authen tication': 'Authentication',
                    'authen tication': 'authentication',
                    'Authoriza tion': 'Authorization',
                    'authoriza tion': 'authorization',
                    'Deplo yment': 'Deployment',
                    'deplo yment': 'deployment',
                    'Documen tation': 'Documentation',
                    'documen tation': 'documentation',
                    'Deliv erables': 'Deliverables',
                    'deliv erables': 'deliverables',
                    'Con ﬁgure': 'Configure',
                    'con ﬁgure': 'configure',
                    'Workﬂow': 'Workflow',
                    'workﬂow': 'workflow',
                    'T esting': 'Testing',
                    't esting': 'testing',
                    'gener ating': 'generating',
                    'r efactoring': 'refactoring',
                    'collabor ator': 'collaborator',
                    'dev eloper': 'developer',
                    'dev elopment': 'development',
                    'elopmen t': 'elopment',
                    'middlew are': 'middleware',
                    'middlew ares': 'middlewares',
                    'con troller': 'controller',
                    'con trollers': 'controllers',
                    'environmen t': 'environment',
                    'notiﬁca tion': 'notification',
                    'notiﬁca tions': 'notifications',
                    'valida tion': 'validation',
                    'integra tion': 'integration',
                    'architec ture': 'architecture',
                    'diagr am': 'diagram',
                    'creden tials': 'credentials',
                    'def ault': 'default',
                    
                    # Words stuck together
                    'BuildaFullstack': 'Build a Fullstack',
                    'needadev': 'need a dev',
                    'needadeveloper': 'need a developer',
                    'whoisane': 'who is an e',
                    'whoisanexpert': 'who is an expert',
                    'shouldbethe': 'should be the',
                    'shouldbe': 'should be',
                    'onlybedoneto': 'only be done to',
                    'onlybedone': 'only be done',
                    'shouldproactively': 'should proactively',
                    'throughoutthepr': 'throughout the pr',
                    'throughoutthe': 'throughout the',
                    'withfrontendandbackend': 'with frontend and backend',
                    'withfrontend': 'with frontend',
                    'andbackend': 'and backend',
                    'GitHubrepo': 'GitHub repo',
                    'Implementroutes': 'Implement routes',
                    'Handletokenstorage': 'Handle token storage',
                    'CreateAddEnquiry': 'Create Add Enquiry',
                    'ReminderforDeveloper': 'Reminder for Developer',
                    'AI-ﬁrst': 'AI-first',
                    'dev elopmenttoolsexpertly': 'development tools expertly',
                    'collabor atorindev': 'collaborator in dev',
                    'indev': 'in dev',
                    'ina': 'in a',
                    'ona': 'on a',
                    'fora': 'for a',
                    'witha': 'with a',
                    'toa': 'to a',
                    'asa': 'as a',
                    
                    # Special characters
                    'ﬂ': 'fl',
                    'ﬁ': 'fi',
                    'ﬀ': 'ff',
                    'ﬃ': 'ffi',
                    'ﬄ': 'ffl',
                }
                
                for old, new in replacements.items():
                    text = text.replace(old, new)
                
                # Step 6: Fix remaining broken words with common suffixes
                text = re.sub(r'([a-z]+)\s+(ment|tion|ing|er|ed|ly|al|ive|ous|ful|ness|able|ible)\b', r'\1\2', text)
                
                # Step 7: Clean multiple spaces
                text = re.sub(r'\s+', ' ', text)
                
                return text.strip() if text.strip() else "No text content found in PDF"
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PDF extraction failed: {str(e)}")

    @staticmethod
    def extract_docx_text(file_path: str) -> str:
        try:
            doc = docx.Document(file_path)
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        text_parts.append(row_text)
            
            text = '\n'.join(text_parts)
            return text.strip() if text.strip() else "No text content found in document"
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"DOCX extraction failed: {str(e)}")

    @staticmethod
    def extract_pptx_text(file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"Slide {slide_num}: " + ' '.join(slide_text))
            
            text = '\n\n'.join(text_parts)
            return text.strip() if text.strip() else "No text content found in presentation"
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPTX extraction failed: {str(e)}")

    @staticmethod
    def extract_text_file(file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Text file reading failed: {str(e)}")

    @staticmethod
    def extract_code_file(file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Code file reading failed: {str(e)}")

# AI Processing
class AIProcessor:
    @staticmethod
    def summarize_text(text: str) -> str:
        """
        INTELLIGENT summarization - Extracts ALL key information comprehensively
        """
        import re
        
        if not text or len(text) < 100:
            return text
        
        # Clean text thoroughly but preserve important content
        text = ' '.join(text.split())
        text = re.sub(r'slide \d+:', '', text, flags=re.IGNORECASE)
        text = re.sub(r'page \d+', '', text, flags=re.IGNORECASE)
        
        # Remove only specific noise patterns, keep project details
        minimal_noise = ['roll no', 'session 20', 'semester']
        for pattern in minimal_noise:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE)
        
        full_text = text.lower()
        
        # Split into sentences - be more inclusive
        sentences = [s.strip() for s in text.split('.') if s.strip()]
        meaningful_sentences = [s for s in sentences if len(s.split()) >= 3]  # Reduced from 5 to 3
        
        if not meaningful_sentences:
            return text[:1000] if len(text) > 1000 else text
        
        # Build comprehensive summary
        summary_parts = []
        
        # 1. IDENTIFY PROJECT/TOPIC
        # Look for project-related keywords in first 20 sentences
        project_indicators = ['project', 'system', 'application', 'platform', 'website', 'portal', 'tool', 'software']
        project_sentence = None
        
        for sent in meaningful_sentences[:20]:
            sent_lower = sent.lower()
            if any(indicator in sent_lower for indicator in project_indicators):
                # This sentence likely describes the project
                project_sentence = sent
                break
        
        if project_sentence:
            # Extract project name - look for capitalized consecutive words
            words = project_sentence.split()
            project_name_parts = []
            for i, word in enumerate(words):
                if word and word[0].isupper() and len(word) > 2 and word.lower() not in ['the', 'and', 'for', 'with', 'using']:
                    project_name_parts.append(word)
                    if len(project_name_parts) >= 6:  # Max 6 words for project name
                        break
                elif project_name_parts:  # Stop if we hit a lowercase word after starting
                    break
            
            if len(project_name_parts) >= 2:
                project_name = ' '.join(project_name_parts)
                summary_parts.append(f"This document describes the {project_name} project.")
            else:
                summary_parts.append(f"This document outlines a software development project.")
        else:
            summary_parts.append("This document contains technical project information.")
        
        # 2. EXTRACT MAIN DESCRIPTION (more comprehensive)
        description_sentences = []
        for sent in meaningful_sentences[:15]:  # Increased from 10 to 15
            sent_clean = sent.strip()
            # Be less restrictive - only skip obvious noise
            if len(sent_clean) > 20 and not any(noise in sent_clean.lower() for noise in ['slide', 'page']):
                description_sentences.append(sent_clean)
                if len(description_sentences) >= 5:  # Increased from 3 to 5
                    break
        
        if description_sentences:
            # Include even more description sentences
            summary_parts.append(' '.join(description_sentences[:6]) + '.')  # Increased from 4 to 6
        
        # 3. TECHNOLOGY STACK (comprehensive detection)
        tech_found = {
            'Frontend': [],
            'Backend': [],
            'Database': [],
            'Tools & Frameworks': []
        }
        
        tech_mapping = {
            # Frontend
            'react': ('Frontend', 'React'), 'angular': ('Frontend', 'Angular'), 'vue': ('Frontend', 'Vue.js'),
            'vite': ('Frontend', 'Vite'), 'next': ('Frontend', 'Next.js'), 'typescript': ('Frontend', 'TypeScript'),
            'javascript': ('Frontend', 'JavaScript'), 'html': ('Frontend', 'HTML'), 'css': ('Frontend', 'CSS'),
            'tailwind': ('Frontend', 'Tailwind CSS'), 'bootstrap': ('Frontend', 'Bootstrap'),
            'material': ('Frontend', 'Material UI'), 'radix': ('Frontend', 'Radix UI'),
            
            # Backend
            'node': ('Backend', 'Node.js'), 'express': ('Backend', 'Express'), 'django': ('Backend', 'Django'),
            'flask': ('Backend', 'Flask'), 'spring': ('Backend', 'Spring Boot'), 'laravel': ('Backend', 'Laravel'),
            'php': ('Backend', 'PHP'), 'python': ('Backend', 'Python'), 'java': ('Backend', 'Java'),
            'ruby': ('Backend', 'Ruby'), 'go': ('Backend', 'Go'), 'rust': ('Backend', 'Rust'),
            
            # Database
            'mongodb': ('Database', 'MongoDB'), 'mysql': ('Database', 'MySQL'), 'postgresql': ('Database', 'PostgreSQL'),
            'postgres': ('Database', 'PostgreSQL'), 'redis': ('Database', 'Redis'), 'sqlite': ('Database', 'SQLite'),
            'oracle': ('Database', 'Oracle'), 'sql': ('Database', 'SQL'), 'nosql': ('Database', 'NoSQL'),
            'firebase': ('Database', 'Firebase'), 'mongoose': ('Database', 'Mongoose'),
            
            # Tools & Frameworks
            'docker': ('Tools & Frameworks', 'Docker'), 'kubernetes': ('Tools & Frameworks', 'Kubernetes'),
            'aws': ('Tools & Frameworks', 'AWS'), 'azure': ('Tools & Frameworks', 'Azure'),
            'git': ('Tools & Frameworks', 'Git'), 'github': ('Tools & Frameworks', 'GitHub'),
            'jenkins': ('Tools & Frameworks', 'Jenkins'), 'nginx': ('Tools & Frameworks', 'Nginx'),
            'apache': ('Tools & Frameworks', 'Apache'), 'jwt': ('Tools & Frameworks', 'JWT'),
            'oauth': ('Tools & Frameworks', 'OAuth'), 'graphql': ('Tools & Frameworks', 'GraphQL'),
            'rest': ('Tools & Frameworks', 'REST API'), 'api': ('Tools & Frameworks', 'API'),
            'jest': ('Tools & Frameworks', 'Jest'), 'mocha': ('Tools & Frameworks', 'Mocha'),
            'eslint': ('Tools & Frameworks', 'ESLint'), 'prettier': ('Tools & Frameworks', 'Prettier'),
            'webpack': ('Tools & Frameworks', 'Webpack'), 'babel': ('Tools & Frameworks', 'Babel')
        }
        
        for keyword, (category, tech_name) in tech_mapping.items():
            if keyword in full_text:
                if tech_name not in tech_found[category]:
                    tech_found[category].append(tech_name)
        
        # Build technology summary
        tech_summary = []
        for category, techs in tech_found.items():
            if techs:
                tech_summary.append(f"{category}: {', '.join(techs)}")
        
        if tech_summary:
            summary_parts.append(f"Technology stack includes {'; '.join(tech_summary)}.")
        
        # 4. KEY FEATURES & MODULES
        features_found = []
        feature_keywords = {
            'Authentication & Authorization': ['authentication', 'login', 'register', 'signup', 'jwt', 'oauth', 'auth', 'authorization'],
            'CRUD Operations': ['crud', 'create', 'read', 'update', 'delete', 'add', 'edit', 'remove'],
            'User Management': ['user management', 'user admin', 'user profile', 'account management'],
            'Dashboard & Analytics': ['dashboard', 'analytics', 'reports', 'statistics', 'metrics', 'visualization'],
            'Search & Filter': ['search', 'filter', 'query', 'find', 'lookup'],
            'API Development': ['api', 'rest', 'restful', 'endpoint', 'web service'],
            'Database Design': ['database', 'schema', 'model', 'entity', 'table', 'collection'],
            'Security': ['security', 'encryption', 'secure', 'protection', 'safety'],
            'Testing': ['testing', 'test', 'unit test', 'integration test', 'qa'],
            'Deployment': ['deployment', 'deploy', 'production', 'hosting', 'server'],
            'Documentation': ['documentation', 'docs', 'readme', 'guide', 'manual'],
            'Responsive Design': ['responsive', 'mobile', 'adaptive', 'device'],
            'Real-time Features': ['real-time', 'realtime', 'live', 'websocket', 'socket'],
            'File Upload': ['upload', 'file upload', 'image upload', 'attachment'],
            'Notification System': ['notification', 'alert', 'email', 'sms', 'push'],
            'Payment Integration': ['payment', 'checkout', 'transaction', 'billing'],
            'Admin Panel': ['admin', 'admin panel', 'administration', 'backend panel']
        }
        
        for feature, keywords in feature_keywords.items():
            if any(kw in full_text for kw in keywords):
                features_found.append(feature)
        
        if features_found:
            if len(features_found) <= 5:
                summary_parts.append(f"Key features include: {', '.join(features_found)}.")
            else:
                summary_parts.append(f"Key features include: {', '.join(features_found[:5])} and {len(features_found) - 5} more modules.")
        
        # 5. DEVELOPMENT METHODOLOGY
        if 'ai' in full_text and any(tool in full_text for tool in ['cursor', 'copilot', 'codeium', 'tabnine', 'ai-first']):
            ai_tools = []
            if 'cursor' in full_text:
                ai_tools.append('Cursor')
            if 'copilot' in full_text:
                ai_tools.append('GitHub Copilot')
            if 'codeium' in full_text:
                ai_tools.append('Codeium')
            if 'tabnine' in full_text:
                ai_tools.append('Tabnine')
            
            if ai_tools:
                summary_parts.append(f"The project follows an AI-first development approach using {', '.join(ai_tools)} for enhanced productivity.")
        
        # 6. PROJECT STRUCTURE
        if 'epic' in full_text:
            epic_keywords = ['project setup', 'authentication', 'user management', 'deployment', 'testing', 'documentation']
            epics_found = [epic for epic in epic_keywords if epic in full_text]
            if epics_found:
                summary_parts.append(f"Development is structured into multiple epics including {', '.join(epics_found)}.")
        
        # 7. DELIVERABLES
        deliverables = []
        if 'source code' in full_text:
            deliverables.append('complete source code')
        if 'docker' in full_text:
            deliverables.append('containerized setup')
        if 'api documentation' in full_text or 'api doc' in full_text:
            deliverables.append('API documentation')
        if 'test suite' in full_text or 'testing' in full_text:
            deliverables.append('test suite')
        if 'deployment' in full_text:
            deliverables.append('deployment configuration')
        if 'documentation' in full_text or 'readme' in full_text:
            deliverables.append('comprehensive documentation')
        
        if deliverables:
            summary_parts.append(f"Project deliverables include: {', '.join(deliverables)}.")
        
        # 8. PROJECT STRUCTURE & IMPLEMENTATION DETAILS
        structure_details = []
        for sent in meaningful_sentences:
            sent_lower = sent.lower()
            # Look for implementation and structure details
            if any(keyword in sent_lower for keyword in ['component', 'module', 'structure', 'architecture', 'implementation', 'design', 'workflow', 'process', 'method', 'approach', 'technique']):
                if len(sent.strip()) > 20:
                    structure_details.append(sent.strip())
                    if len(structure_details) >= 3:
                        break
        
        if structure_details:
            summary_parts.append(f"Implementation approach: {' '.join(structure_details[:2])}")
        
        # 9. ADDITIONAL CONTEXT (from any remaining meaningful content)
        if len(meaningful_sentences) > 10:
            context_sentences = []
            for sent in meaningful_sentences[10:]:
                if len(sent.strip()) > 25 and len(context_sentences) < 3:
                    context_sentences.append(sent.strip())
            
            if context_sentences:
                summary_parts.append(f"Additional context: {' '.join(context_sentences[:2])}")
        
        # 8. ENSURE PROPORTIONAL LENGTH BASED ON CONTENT
        result = ' '.join(summary_parts)
        
        # Calculate proportional summary length based on original content
        original_length = len(text)
        target_summary_length = 0
        
        if original_length < 1000:
            target_summary_length = 400  # Increased from 300
        elif original_length < 3000:
            target_summary_length = 700  # Increased from 500
        elif original_length < 6000:
            target_summary_length = 1000  # Increased from 800
        else:
            target_summary_length = 1300  # Increased from 1000
        
        # If summary is too short compared to content, add more details
        if len(result) < target_summary_length and len(meaningful_sentences) > 5:
            # Add more context from meaningful sentences
            additional_sentences = []
            for sent in meaningful_sentences[5:15]:  # Take sentences 6-15
                sent_clean = sent.strip()
                if len(sent_clean) > 15:  # Reduced threshold from 30 to 15
                    additional_sentences.append(sent_clean)
                    if len(' '.join(additional_sentences)) > (target_summary_length - len(result)):
                        break
            
            if additional_sentences:
                result += ' ' + ' '.join(additional_sentences)
        
        # Add more technical details if content is very long
        if original_length > 4000 and len(result) < target_summary_length:
            # Look for more specific technical details
            technical_sentences = []
            for sent in meaningful_sentences:
                sent_lower = sent.lower()
                # Look for sentences with technical keywords
                if any(tech in sent_lower for tech in ['framework', 'library', 'database', 'server', 'client', 'architecture', 'design', 'implementation', 'feature', 'module', 'component']):
                    if len(sent.strip()) > 40:
                        technical_sentences.append(sent.strip())
                        if len(technical_sentences) >= 3:
                            break
            
            if technical_sentences:
                result += ' Additional technical details: ' + ' '.join(technical_sentences[:2])
        
        # Clean up
        result = result.replace('..', '.').replace('  ', ' ')
        result = re.sub(r'\s+', ' ', result).strip()
        
        # Final length adjustment - ensure we meet the target (more aggressive)
        if len(result) < target_summary_length * 0.7 and len(meaningful_sentences) > 8:  # Reduced from 0.8 to 0.7
            # Add more sentences to reach target length
            extra_sentences = meaningful_sentences[8:20]  # Increased from 12 to 20
            extra_content = ' '.join([s.strip() for s in extra_sentences if len(s.strip()) > 10])  # Reduced from 25 to 10
            if extra_content:
                result += ' Additional details: ' + extra_content[:target_summary_length - len(result)]
        
        # Cap maximum length
        if len(result) > target_summary_length * 1.2:
            result = result[:int(target_summary_length * 1.2)]
            # Ensure it ends with a complete sentence
            last_period = result.rfind('.')
            if last_period > len(result) * 0.8:
                result = result[:last_period + 1]
        
        return result
    
    @staticmethod
    def explain_code(code: str, file_extension: str) -> str:
        """
        🚀 ADVANCED UNIVERSAL CODE ANALYZER
        Uses the new advanced analyzer that understands ALL programming languages
        Supports: Java, Python, JavaScript, C++, C, C#, PHP, Ruby, Go, Rust, TypeScript, etc.
        """
        # Use the advanced universal analyzer
        return analyze_code_advanced_universal(code, file_extension)

# Audio generation
class AudioGenerator:
    @staticmethod
    def clean_text_for_audio(text: str) -> str:
        """
        Clean text for better audio generation
        Removes visual formatting elements that don't sound good when read aloud
        """
        import re
        
        # Remove long lines of equal signs, dashes, or other repeated characters
        text = re.sub(r'[=\-_*#]{10,}', '', text)
        
        # Remove emoji and special symbols that don't read well
        text = re.sub(r'[🎯🔧💡📊⚡🚀✅❌🎉🎵📥📱🌟🔍🎨🎓📚🔮]', '', text)
        
        # Remove bullet points and replace with "Point:"
        text = re.sub(r'[•▶✓]', 'Point:', text)
        
        # Clean up multiple spaces and newlines
        text = re.sub(r'\n{3,}', '\n\n', text)  # Max 2 newlines
        text = re.sub(r' {2,}', ' ', text)      # Max 1 space
        
        # Remove standalone numbers at start of lines (like "1. 2. 3.")
        text = re.sub(r'^\d+\.\s*', '', text, flags=re.MULTILINE)
        
        # Replace common symbols with words
        replacements = {
            '&': ' and ',
            '@': ' at ',
            '#': ' hash ',
            '%': ' percent ',
            '+': ' plus ',
            '=': ' equals ',
            '<': ' less than ',
            '>': ' greater than ',
            '|': ' or ',
            '()': '',
            '[]': '',
            '{}': '',
            '→': ' leads to ',
            '←': ' comes from ',
            '↑': ' up ',
            '↓': ' down '
        }
        
        for symbol, word in replacements.items():
            text = text.replace(symbol, word)
        
        # Clean up extra whitespace
        text = ' '.join(text.split())
        
        return text.strip()

    @staticmethod
    def generate_audio(text: str, language: str = 'en') -> str:
        try:
            # Clean text for better audio experience
            cleaned_text = AudioGenerator.clean_text_for_audio(text)
            
            # Generate unique filename
            audio_id = str(uuid.uuid4())
            audio_path = f"audio/{audio_id}.mp3"
            
            # Create audio directory if it doesn't exist
            os.makedirs("audio", exist_ok=True)
            
            # Generate audio using gTTS with cleaned text
            tts = gTTS(text=cleaned_text, lang=language, slow=False)
            tts.save(audio_path)
            
            return audio_path
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Audio generation failed: {str(e)}")

# API Endpoints
@app.get("/")
async def root():
    return {"message": "🤖 AI Document Reader Service is running!", "status": "active"}

@app.post("/process-file")
async def process_file(request: FileProcessRequest):
    try:
        file_path = request.filePath
        file_type = request.fileType.lower()
        
        print(f"Processing file: {file_path}")
        print(f"File type: {file_type}")
        
        # Extract content based on file type
        if file_type == '.pdf':
            content = DocumentProcessor.extract_pdf_text(file_path)
            content_type = 'document'
        elif file_type in ['.docx', '.doc']:
            content = DocumentProcessor.extract_docx_text(file_path)
            content_type = 'document'
        elif file_type == '.pptx':
            content = DocumentProcessor.extract_pptx_text(file_path)
            content_type = 'document'
        elif file_type == '.txt':
            content = DocumentProcessor.extract_text_file(file_path)
            content_type = 'document'
        elif file_type in ['.py', '.js', '.java', '.cpp', '.c']:
            content = DocumentProcessor.extract_code_file(file_path)
            content_type = 'code'
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        print(f"Successfully extracted {len(content)} characters")
        
        return {
            "success": True,
            "contentType": content_type,
            "fileType": file_type,
            "contentLength": len(content),
            "preview": content[:200] + "..." if len(content) > 200 else content
        }
        
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/get-content")
async def get_content(request: ContentRequest):
    try:
        file_path = request.filePath
        content_type = request.contentType
        
        print(f"Getting content for: {file_path}, type: {content_type}")
        
        # Determine file type
        file_extension = Path(file_path).suffix.lower()
        print(f"File extension: {file_extension}")
        
        # Extract content
        if file_extension == '.pdf':
            full_content = DocumentProcessor.extract_pdf_text(file_path)
        elif file_extension in ['.docx', '.doc']:
            full_content = DocumentProcessor.extract_docx_text(file_path)
        elif file_extension == '.pptx':
            full_content = DocumentProcessor.extract_pptx_text(file_path)
        elif file_extension == '.txt':
            full_content = DocumentProcessor.extract_text_file(file_path)
        elif file_extension in ['.py', '.js', '.java', '.cpp', '.c']:
            full_content = DocumentProcessor.extract_code_file(file_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        print(f"Extracted {len(full_content)} characters")
        
        # Process content based on request type
        if file_extension in ['.py', '.js', '.java', '.cpp', '.c']:
            # For code files, always explain the code
            print("Processing as code file")
            processed_content = AIProcessor.explain_code(full_content, file_extension)
        else:
            # For documents
            if content_type == 'summary':
                print("Generating summary")
                processed_content = AIProcessor.summarize_text(full_content)
            else:
                print("Returning full content")
                processed_content = full_content
        
        print(f"Processed content length: {len(processed_content)}")
        
        return {
            "success": True,
            "content": processed_content,
            "contentType": content_type,
            "isCode": file_extension in ['.py', '.js', '.java', '.cpp', '.c']
        }
        
    except Exception as e:
        print(f"ERROR in get_content: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate-audio")
async def generate_audio(request: AudioRequest):
    try:
        audio_path = AudioGenerator.generate_audio(request.text, request.language)
        
        return {
            "success": True,
            "audioPath": audio_path,
            "audioUrl": f"http://localhost:8001/{audio_path}"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Serve audio files
from fastapi.staticfiles import StaticFiles
app.mount("/audio", StaticFiles(directory="audio"), name="audio")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
